import argparse
import re
import sys
from pathlib import Path
from pptx import Presentation

# デフォルト設定
DEFAULT_TEMPLATE = "template.pptx"
DEFAULT_OUTPUT = "output.pptx"
DEFAULT_MARKDOWN = "slides.md"

LAYOUT_INDEX = 1
PH_TITLE_IDX = 0
PH_BODY1_IDX = 1
PH_BODY2_IDX = 3


def parse_markdown(md_file_path: str):
    """
    Markdownファイルを '---' で区切り、各スライドの「タイトル」「本文1」「本文2」を抽出します。
    """
    with open(md_file_path, "r", encoding="utf-8") as f:
        content = f.read()

    # スライド単位で分割 ('---' で区切られている前提)
    raw_slides = re.split(r"\n---\n", content)
    slides_data = []

    for raw_slide in raw_slides:
        slide_info = {"title": "", "body1": "", "body2": ""}

        # タイトル (# タイトル)
        title_match = re.search(r"^#\s+(.+)$", raw_slide, re.MULTILINE)
        if title_match:
            slide_info["title"] = title_match.group(1).strip()

        # ## 見出しブロックを順番に取得 (1番目を本文1、2番目を本文2)
        sections = re.findall(
            r"^##\s*.*$\n([\s\S]*?)(?=(?:^##|\Z))", raw_slide, re.MULTILINE
        )
        if len(sections) >= 1:
            slide_info["body1"] = sections[0].strip()
        if len(sections) >= 2:
            slide_info["body2"] = sections[1].strip()

        slides_data.append(slide_info)

    return slides_data


import copy


def generate_pptx(slides_data, template_path: str, output_path: str, layout_index: int = LAYOUT_INDEX):
    """Markdownから取得したデータをテンプレートのPPTXに流し込みます。"""
    prs = Presentation(template_path)
    slide_layout = prs.slide_layouts[layout_index]

    for slide_data in slides_data:
        # スライド追加
        slide = prs.slides.add_slide(slide_layout)

        # python-pptxでadd_slide時に引き継がれないレイアウト上のプレースホルダー（ページ番号、フッター、日付など）を保持
        existing_ph_idxs = {ph.placeholder_format.idx for ph in slide.placeholders}
        for layout_ph in slide_layout.placeholders:
            if layout_ph.placeholder_format.idx not in existing_ph_idxs:
                slide.shapes._spTree.append(copy.deepcopy(layout_ph._element))

        placeholders = list(slide.placeholders)

        # タイトル設定
        if len(placeholders) > PH_TITLE_IDX and placeholders[PH_TITLE_IDX].has_text_frame:
            placeholders[PH_TITLE_IDX].text = slide_data["title"]

        # 本文1設定
        if len(placeholders) > PH_BODY1_IDX and placeholders[PH_BODY1_IDX].has_text_frame:
            placeholders[PH_BODY1_IDX].text = slide_data["body1"]

        # 本文2設定
        if len(placeholders) > PH_BODY2_IDX and placeholders[PH_BODY2_IDX].has_text_frame:
            placeholders[PH_BODY2_IDX].text = slide_data["body2"]

    prs.save(output_path)
    print(f"作成完了: {output_path} ({len(slides_data)}枚のスライド)")


def main():
    parser = argparse.ArgumentParser(
        description="MarkdownからPowerPoint（.pptx）スライドを自動生成するツール"
    )
    parser.add_argument(
        "-i",
        "--input",
        dest="input",
        default=DEFAULT_MARKDOWN,
        help=f"入力Markdownファイルパス (デフォルト: {DEFAULT_MARKDOWN})",
    )
    parser.add_argument(
        "-t",
        "--template",
        default=DEFAULT_TEMPLATE,
        help=f"PowerPointテンプレートファイルパス (デフォルト: {DEFAULT_TEMPLATE})",
    )
    parser.add_argument(
        "-o",
        "--output",
        default=DEFAULT_OUTPUT,
        help=f"出力PowerPointファイルパス (デフォルト: {DEFAULT_OUTPUT})",
    )
    parser.add_argument(
        "--layout",
        type=int,
        default=LAYOUT_INDEX,
        help=f"スライドレイアウト番号 (デフォルト: {LAYOUT_INDEX})",
    )

    args = parser.parse_args()

    if not Path(args.input).exists():
        print(f"エラー: Markdownファイルが見つかりません: {args.input}", file=sys.stderr)
        sys.exit(1)

    if not Path(args.template).exists():
        print(f"エラー: テンプレートファイルが見つかりません: {args.template}", file=sys.stderr)
        sys.exit(1)

    data = parse_markdown(args.input)
    generate_pptx(data, args.template, args.output, layout_index=args.layout)


if __name__ == "__main__":
    main()
